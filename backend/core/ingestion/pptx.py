from pptx import Presentation

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extracts text from PPTX slides and speaker notes.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i+1} ---\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                    
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    text += f"Notes: {notes}\n"
                    
            text += "\n"
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
    return text
